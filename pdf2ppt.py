import argparse
import fitz  # PyMuPDF
from pdf2image import convert_from_path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
import os

def image_based_conversion(pdf_path, pptx_path):
    # Convert PDF pages to images
    images = convert_from_path(pdf_path)
    
    # Create a PowerPoint presentation
    prs = Presentation()
    prs.slide_width = Inches(11.7)
    prs.slide_height = Inches(8.3)
    
    for img in images:
        slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
        img_path = 'temp_image.png'
        img.save(img_path, 'PNG')
        slide.shapes.add_picture(img_path, Inches(0), Inches(0), width=prs.slide_width, height=prs.slide_height)
        os.remove(img_path)  # Clean up temp image file
    
    prs.save(pptx_path)
    print(f"Image-based PowerPoint saved as {pptx_path}")

def text_based_conversion(pdf_path, pptx_path):
    # Open PDF with PyMuPDF
    doc = fitz.open(pdf_path)
    prs = Presentation()
    
    for page_num in range(doc.page_count):
        page = doc.load_page(page_num)
        text = page.get_text("text")
        
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8.5), Inches(5.5))
        text_frame = textbox.text_frame
        p = text_frame.add_paragraph()
        p.text = text
        p.font.size = Pt(12)
        text_frame.word_wrap = True
        p.alignment = PP_ALIGN.LEFT
    
    prs.save(pptx_path)
    print(f"Text-based PowerPoint saved as {pptx_path}")

def main():
    parser = argparse.ArgumentParser(description="Convert PDF to PowerPoint (PPTX).")
    parser.add_argument("pdf_file", help="Path to the PDF file to convert.")
    parser.add_argument("--output", "-o", default="output_presentation.pptx", help="Output PPTX filename.")
    parser.add_argument("--mode", choices=["image", "text"], default="image",
                        help="Conversion mode: 'image' for page images or 'text' for editable text.")
    
    args = parser.parse_args()

    if args.mode == "image":
        image_based_conversion(args.pdf_file, args.output)
    elif args.mode == "text":
        text_based_conversion(args.pdf_file, args.output)
    else:
        print("Invalid mode selected. Use 'image' or 'text'.")

if __name__ == "__main__":
    main()

